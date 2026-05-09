#!/usr/bin/env python3
"""Generate teaching PPT, Excel dashboard, Gantt image, HTML dashboard, and brief.

This script is intentionally simple and deterministic so it can be used in a teaching demo.
"""
from __future__ import annotations

import argparse
import json
from datetime import datetime
from pathlib import Path
import yaml

from pptx import Presentation
from pptx.util import Inches, Pt
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.chart import BarChart, Reference
import matplotlib.pyplot as plt

STATUS_LABEL = {
    "done": "Done",
    "in_progress": "In progress",
    "not_started": "Not started",
    "needs_review": "Needs review",
}


def load_request(path: Path) -> str:
    return path.read_text(encoding="utf-8").strip()


def derive_status(section_id: str, progress: dict) -> str:
    if section_id == progress.get("active_demo"):
        return "in_progress"
    current_min = progress.get("current_minute", 0)
    return "done" if section_id in ["demo_00"] and current_min >= 10 else "not_started"


def build_schedule(template: dict, progress: dict):
    rows = []
    for sec in template["sections"]:
        status = derive_status(sec["id"], progress)
        rows.append({**sec, "status": status, "owner": progress.get("owners", {}).get(sec["id"], "TBD")})
    return rows


def write_brief(output_dir: Path, request: str, template: dict, progress: dict, schedule: list[dict]):
    completed = ", ".join(progress.get("completed", [])) or "暂无"
    questions = "\n".join(f"- {q}" for q in progress.get("user_questions", [])) or "- 暂无"
    schedule_md = "\n".join(
        f"- {r['start_min']:02d}-{r['end_min']:02d} min | {r['title']} | {STATUS_LABEL[r['status']]}"
        for r in schedule
    )
    review = "\n".join(f"- [ ] {q}" for q in template.get("review_questions", []))
    text = f"""# 教学进度简报

## 用户一句话需求

{request}

## 当前进展

- 当前阶段：{progress.get('current_stage', 'TBD')}
- 已完成：{completed}
- 当前分钟：{progress.get('current_minute', 0)} / {template.get('duration_minutes', 60)}

## 总体时间规划

{schedule_md}

## 用户问题

{questions}

## 人工 Review Checklist

{review}

## 讲师提示

这份材料由规则化脚本生成，适合课堂实时更新。正式对外发布前需要人工确认时间、问题表述和最终措辞。
"""
    (output_dir / "teaching_brief.md").write_text(text, encoding="utf-8")


def write_pptx(output_dir: Path, request: str, template: dict, progress: dict, schedule: list[dict]):
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = template["course_title"]
    title_slide.placeholders[1].text = "规则化生成：PPT + Excel + 甘特图 + Dashboard"

    def bullet_slide(title: str, bullets: list[str]):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(20)
        return slide

    bullet_slide("用户一句话需求", [request])
    bullet_slide("当前进展", [
        f"当前阶段：{progress.get('current_stage')}",
        f"当前分钟：{progress.get('current_minute')} / {template.get('duration_minutes')}",
        "已完成：" + ", ".join(progress.get("completed", [])),
    ])
    bullet_slide("60 分钟规划", [
        f"{r['start_min']}-{r['end_min']} min: {r['title']} ({STATUS_LABEL[r['status']]})" for r in schedule
    ])
    bullet_slide("用户问题", progress.get("user_questions", ["暂无"])[0:5])
    bullet_slide("Takeaway", [
        "Prompt 解决一次，Skill 解决一类。",
        "模板和脚本让生成结果稳定。",
        "权限和人工 review 决定安全边界。",
        "OpenCode 适合把规则、命令、工具和 Skill 串起来。",
    ])
    prs.save(output_dir / "teaching_progress_deck.pptx")


def write_gantt_png(output_dir: Path, schedule: list[dict], duration: int):
    labels = [r["id"] for r in schedule]
    starts = [r["start_min"] for r in schedule]
    widths = [r["end_min"] - r["start_min"] for r in schedule]
    y = range(len(schedule))
    plt.figure(figsize=(10, 4.5))
    plt.barh(y, widths, left=starts)
    plt.yticks(y, labels, fontsize=8)
    plt.xlabel("Minute")
    plt.xlim(0, duration)
    plt.title("60-minute teaching Gantt")
    plt.tight_layout()
    plt.savefig(output_dir / "teaching_gantt.png", dpi=160)
    plt.close()


def write_xlsx(output_dir: Path, template: dict, progress: dict, schedule: list[dict]):
    wb = Workbook()
    ws = wb.active
    ws.title = "Dashboard"
    ws["A1"] = template["course_title"]
    ws["A1"].font = Font(size=16, bold=True)
    ws.merge_cells("A1:H1")
    ws["A3"] = "KPI"
    ws["A3"].font = Font(bold=True)
    kpis = [
        ("Current minute", progress.get("current_minute", 0)),
        ("Remaining minute", template.get("duration_minutes", 60) - progress.get("current_minute", 0)),
        ("Active stage", progress.get("current_stage", "TBD")),
        ("Open questions", len(progress.get("user_questions", []))),
    ]
    for idx, (k, v) in enumerate(kpis, start=4):
        ws[f"A{idx}"] = k
        ws[f"B{idx}"] = v
    headers = ["ID", "Title", "Start", "End", "Duration", "Status", "Owner", "Objective"]
    start_row = 10
    for col, h in enumerate(headers, start=1):
        c = ws.cell(row=start_row, column=col, value=h)
        c.font = Font(bold=True, color="FFFFFF")
        c.fill = PatternFill("solid", fgColor="1F4E78")
        c.alignment = Alignment(horizontal="center")
    for r_idx, row in enumerate(schedule, start=start_row + 1):
        values = [row["id"], row["title"], row["start_min"], row["end_min"], row["end_min"] - row["start_min"], STATUS_LABEL[row["status"]], row["owner"], row["objective"]]
        for c_idx, val in enumerate(values, start=1):
            ws.cell(row=r_idx, column=c_idx, value=val)
    dv = DataValidation(type="list", formula1='"Done,In progress,Not started,Needs review"', allow_blank=False)
    ws.add_data_validation(dv)
    dv.add(f"F{start_row+1}:F{start_row+len(schedule)}")
    chart = BarChart()
    chart.title = "Section Duration"
    chart.y_axis.title = "Minutes"
    chart.x_axis.title = "Section"
    data = Reference(ws, min_col=5, min_row=start_row, max_row=start_row + len(schedule))
    cats = Reference(ws, min_col=2, min_row=start_row + 1, max_row=start_row + len(schedule))
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(cats)
    ws.add_chart(chart, "J3")
    ws2 = wb.create_sheet("Questions")
    ws2.append(["Question", "Owner", "Status", "Suggested handling"])
    for q in progress.get("user_questions", []):
        ws2.append([q, "讲师", "Open", "放入 Q&A 或当前 demo 中回答"])
    ws3 = wb.create_sheet("Manifest")
    ws3.append(["Generated at", datetime.now().isoformat(timespec="seconds")])
    ws3.append(["Generator", "scripts/generate_training_artifacts.py"])
    ws3.append(["Human review required", "Yes"])
    for sheet in wb.worksheets:
        for col in range(1, min(sheet.max_column, 12) + 1):
            sheet.column_dimensions[chr(64 + col)].width = 18
    wb.save(output_dir / "teaching_dashboard.xlsx")


def write_html(output_dir: Path, progress: dict, schedule: list[dict]):
    cards = "".join(f"<tr><td>{r['title']}</td><td>{r['start_min']}-{r['end_min']}</td><td>{STATUS_LABEL[r['status']]}</td><td>{r['objective']}</td></tr>" for r in schedule)
    questions = "".join(f"<li>{q}</li>" for q in progress.get("user_questions", []))
    html = f"""<!doctype html><html><head><meta charset='utf-8'><title>Teaching Dashboard</title>
<style>body{{font-family:Arial,sans-serif;margin:32px}} table{{border-collapse:collapse;width:100%}} td,th{{border:1px solid #ddd;padding:8px}} th{{background:#eee}}</style></head>
<body><h1>Teaching Dashboard</h1><p><b>Current stage:</b> {progress.get('current_stage')}</p><p><b>Current minute:</b> {progress.get('current_minute')}</p>
<h2>Schedule</h2><table><tr><th>Section</th><th>Time</th><th>Status</th><th>Objective</th></tr>{cards}</table>
<h2>User Questions</h2><ul>{questions}</ul><p><i>Generated for teaching demo. Human review required before sharing.</i></p></body></html>"""
    (output_dir / "dashboard.html").write_text(html, encoding="utf-8")


def write_manifest(output_dir: Path, request: str, template: dict, progress: dict, schedule: list[dict]):
    manifest = {
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "request_chars": len(request),
        "outputs": ["teaching_brief.md", "teaching_progress_deck.pptx", "teaching_dashboard.xlsx", "teaching_gantt.png", "dashboard.html"],
        "human_review_required": True,
        "review_reasons": template.get("review_questions", []),
        "schedule_count": len(schedule),
        "current_stage": progress.get("current_stage"),
    }
    (output_dir / "generation_manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--request", required=True)
    ap.add_argument("--progress", required=True)
    ap.add_argument("--template", required=True)
    ap.add_argument("--output-dir", required=True)
    args = ap.parse_args()
    request = load_request(Path(args.request))
    progress = json.loads(Path(args.progress).read_text(encoding="utf-8"))
    template = yaml.safe_load(Path(args.template).read_text(encoding="utf-8"))
    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    schedule = build_schedule(template, progress)
    write_brief(output_dir, request, template, progress, schedule)
    write_pptx(output_dir, request, template, progress, schedule)
    write_gantt_png(output_dir, schedule, template.get("duration_minutes", 60))
    write_xlsx(output_dir, template, progress, schedule)
    write_html(output_dir, progress, schedule)
    write_manifest(output_dir, request, template, progress, schedule)
    print(json.dumps({"ok": True, "output_dir": str(output_dir), "files": sorted(p.name for p in output_dir.iterdir())}, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
