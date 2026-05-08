#!/usr/bin/env python3
from __future__ import annotations

import json
from datetime import datetime
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]


def load_json(path: Path):
    with path.open(encoding='utf-8') as f:
        return json.load(f)


def render_workstream_table(workstreams):
    lines = ['| Workstream | Status | Progress | Evidence | Next | Risk |', '|---|---:|---:|---|---|---|']
    for item in workstreams:
        evidence = '<br>'.join(item.get('evidence', []))
        lines.append(f"| {item['name']} | {item['status']} | {item['progress']}% | {evidence} | {item['next']} | {item['risk']} |")
    return '\n'.join(lines)


def render_bullets(items):
    return '\n'.join(f'- {item}' for item in items) if items else '- None reported.'


def build_markdown(progress, rules, template):
    next_steps = [f"{w['name']}: {w['next']}" for w in progress['workstreams']]
    checklist = rules['human_review_questions']
    rendered = template
    rendered = rendered.replace('{{project_name}}', progress['project_name'])
    rendered = rendered.replace('{{period}}', progress['period'])
    rendered = rendered.replace('{{summary}}', progress['summary'])
    rendered = rendered.replace('{{workstream_table}}', render_workstream_table(progress['workstreams']))
    rendered = rendered.replace('{{blockers}}', render_bullets(progress.get('blockers', [])))
    rendered = rendered.replace('{{next_steps}}', render_bullets(next_steps))
    rendered = rendered.replace('{{review_checklist}}', render_bullets(checklist))
    return rendered


def build_pptx(progress, rules, output_path: Path):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception as exc:  # pragma: no cover - dependency fallback
        return {'created': False, 'reason': f'python-pptx unavailable: {exc}'}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title(slide, title, subtitle=''):
        shape = slide.shapes.title
        shape.text = title
        shape.text_frame.paragraphs[0].font.size = Pt(30)
        if subtitle:
            box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.8), Inches(1.0))
            box.text_frame.text = subtitle

    def add_bullets(slide, title, bullets):
        add_title(slide, title)
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.8), Inches(5.3))
        tf = box.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    add_title(slide, progress['project_name'], f"{progress['period']} - {progress['report_date']} - owner: {progress['owner']}")

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_bullets(slide, 'Executive summary', [progress['summary'], f"Demos included: {progress['metrics']['demo_count']}", 'Validation is scripted; final messaging requires human review.'])

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_bullets(slide, 'Progress by workstream', [f"{w['name']}: {w['status']} ({w['progress']}%), risk={w['risk']}" for w in progress['workstreams']])

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_bullets(slide, 'Risks and blockers', progress.get('blockers', []) + progress.get('decisions_needed', []))

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_bullets(slide, 'Human review checklist', rules['human_review_questions'])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return {'created': True, 'slide_count': len(prs.slides)}


def main() -> int:
    progress = load_json(ROOT / 'inputs/progress_update.json')
    rules = load_json(ROOT / 'rules/report_rules.json')
    template = (ROOT / 'templates/project_status_template.md').read_text(encoding='utf-8')
    output = ROOT / 'output'
    output.mkdir(exist_ok=True)

    md = build_markdown(progress, rules, template)
    md_path = output / 'project_status_report.md'
    md_path.write_text(md, encoding='utf-8')

    pptx_info = build_pptx(progress, rules, output / 'project_status_report.pptx')

    manifest = {
        'generated_at': datetime.utcnow().isoformat() + 'Z',
        'inputs': ['inputs/progress_update.json', 'rules/report_rules.json', 'templates/project_status_template.md'],
        'outputs': ['output/project_status_report.md'] + (['output/project_status_report.pptx'] if pptx_info.get('created') else []),
        'pptx': pptx_info,
        'human_review_required': True
    }
    (output / 'report_manifest.json').write_text(json.dumps(manifest, indent=2), encoding='utf-8')
    print(f'Generated {md_path}')
    return 0


if __name__ == '__main__':
    raise SystemExit(main())
