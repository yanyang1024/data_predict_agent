#!/usr/bin/env python3
"""生成课程进度 PPT、最小 XLSX 看板和 Gantt HTML。

设计意图：
- LLM/Agent 负责把用户一句话归一化为 inputs/course_status.json；
- 脚本只负责稳定渲染，不负责判断课程内容是否真实正确。
"""
from __future__ import annotations
import argparse, json, html, zipfile
from pathlib import Path
from datetime import datetime
from xml.sax.saxutils import escape


def col_name(n: int) -> str:
    s = ""
    while n:
        n, r = divmod(n - 1, 26)
        s = chr(65 + r) + s
    return s


def cell_xml(row: int, col: int, value):
    ref = f"{col_name(col)}{row}"
    if isinstance(value, (int, float)) and not isinstance(value, bool):
        return f'<c r="{ref}"><v>{value}</v></c>'
    value = "" if value is None else str(value)
    return f'<c r="{ref}" t="inlineStr"><is><t>{escape(value)}</t></is></c>'


def sheet_xml(rows):
    xml_rows = []
    for r_idx, row in enumerate(rows, start=1):
        cells = ''.join(cell_xml(r_idx, c_idx, v) for c_idx, v in enumerate(row, start=1))
        xml_rows.append(f'<row r="{r_idx}">{cells}</row>')
    return '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>' + \
        '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><sheetData>' + ''.join(xml_rows) + '</sheetData></worksheet>'


def write_xlsx(path: Path, sheets: dict[str, list[list[object]]]):
    content_types = ['<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">',
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>',
        '<Default Extension="xml" ContentType="application/xml"/>',
        '<Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>']
    for i in range(1, len(sheets)+1):
        content_types.append(f'<Override PartName="/xl/worksheets/sheet{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>')
    content_types.append('</Types>')

    workbook_sheets = []
    rels = ['<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">']
    for i, name in enumerate(sheets, start=1):
        workbook_sheets.append(f'<sheet name="{escape(name)}" sheetId="{i}" r:id="rId{i}"/>')
        rels.append(f'<Relationship Id="rId{i}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet{i}.xml"/>')
    rels.append('</Relationships>')

    workbook_xml = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>' + \
        '<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><sheets>' + ''.join(workbook_sheets) + '</sheets></workbook>'
    root_rels = '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>' + \
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/></Relationships>'

    with zipfile.ZipFile(path, 'w', zipfile.ZIP_DEFLATED) as z:
        z.writestr('[Content_Types].xml', ''.join(content_types))
        z.writestr('_rels/.rels', root_rels)
        z.writestr('xl/workbook.xml', workbook_xml)
        z.writestr('xl/_rels/workbook.xml.rels', ''.join(rels))
        for i, rows in enumerate(sheets.values(), start=1):
            z.writestr(f'xl/worksheets/sheet{i}.xml', sheet_xml(rows))


def build_pptx(status: dict, request_text: str, path: Path):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception:
        fallback = path.with_suffix('.pptx.md')
        fallback.write_text('# PPT 降级输出\n\n本环境缺少 python-pptx，已生成 Markdown 摘要。\n', encoding='utf-8')
        return str(fallback)

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = status['course_title']
    slide.placeholders[1].text = f"当前阶段：{status['current_phase']}\n总时长：{status['total_minutes']} 分钟"

    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = '1 小时时间规划'
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for d in status['demos']:
        p = tf.add_paragraph()
        p.text = f"{d['minutes']}  {d['id']} {d['name']}：{d['key_message']}"
        p.level = 0

    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = '四个 Demo 当前进展'
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for d in status['demos']:
        p = tf.add_paragraph()
        p.text = f"{d['id']}｜{d['status']}｜{d['progress']}%｜{d['name']}"
        p.level = 0

    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = '用户问题与风险'
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.add_paragraph().text = '用户问题：'
    for q in status['user_questions']:
        p = tf.add_paragraph(); p.text = q; p.level = 1
    tf.add_paragraph().text = '风险：'
    for r in status['risks']:
        p = tf.add_paragraph(); p.text = r; p.level = 1

    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = '上下文沉淀结构'
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for item in ['AGENTS.md：长期规则', 'Command：高频入口', 'Skill：流程和 Stop Rules', 'references/templates：业务知识和模板', 'scripts/tools/API：受控执行', 'validator：自动检查', 'Human Review：逻辑确认']:
        p = tf.add_paragraph(); p.text = item

    slide = prs.slides.add_slide(content_layout)
    slide.shapes.title.text = '下一步行动'
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for a in status['next_actions']:
        p = tf.add_paragraph(); p.text = a
    tf.add_paragraph().text = '人工确认：请讲师确认进度、问题、风险和讲解顺序是否准确。'

    prs.save(path)
    return str(path)


def build_html_gantt(status: dict, path: Path):
    bars = []
    for d in status['demos']:
        width = max(2, min(100, int(d['progress'])))
        bars.append(f"""
        <tr>
          <td>{html.escape(d['id'])}</td><td>{html.escape(d['name'])}</td><td>{html.escape(d['minutes'])}</td><td>{html.escape(d['status'])}</td>
          <td><div class='bar-bg'><div class='bar' style='width:{width}%'>{width}%</div></div></td>
        </tr>""")
    doc = f"""<!doctype html><html lang='zh-CN'><head><meta charset='utf-8'>
<title>课程 Gantt Dashboard</title>
<style>
body {{ font-family: Arial, 'Microsoft YaHei', sans-serif; margin: 32px; color:#172033; }}
table {{ border-collapse: collapse; width: 100%; }} th, td {{ border:1px solid #d0d7de; padding:8px; }} th {{ background:#f6f8fa; }}
.bar-bg {{ background:#eef2f7; height:22px; border-radius:11px; overflow:hidden; }} .bar {{ background:#4f46e5; color:white; text-align:center; height:22px; line-height:22px; font-size:12px; }}
.card {{ padding:16px; background:#f8fafc; border:1px solid #e5e7eb; border-radius:10px; margin: 12px 0; }}
</style></head><body>
<h1>{html.escape(status['course_title'])}</h1>
<div class='card'>当前阶段：{html.escape(status['current_phase'])}｜总体进度：{status['overall_progress_percent']}%｜总时长：{status['total_minutes']} 分钟</div>
<table><thead><tr><th>ID</th><th>Demo</th><th>时间</th><th>状态</th><th>进度</th></tr></thead><tbody>{''.join(bars)}</tbody></table>
<h2>用户问题</h2><ul>{''.join(f'<li>{html.escape(q)}</li>' for q in status['user_questions'])}</ul>
<h2>风险与人工确认</h2><ul>{''.join(f'<li>{html.escape(r)}</li>' for r in status['risks'])}</ul>
<p><strong>提示：</strong>此看板由脚本生成，课程内容准确性仍需讲师确认。</p>
</body></html>"""
    path.write_text(doc, encoding='utf-8')


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--status', required=True)
    ap.add_argument('--request', required=True)
    ap.add_argument('--output', required=True)
    args = ap.parse_args()
    out = Path(args.output); out.mkdir(parents=True, exist_ok=True)
    status = json.loads(Path(args.status).read_text(encoding='utf-8'))
    request_text = Path(args.request).read_text(encoding='utf-8')

    ppt_path = out / 'course_update.pptx'
    actual_ppt = build_pptx(status, request_text, ppt_path)

    dashboard_rows = [['课程', status['course_title']], ['总时长', status['total_minutes']], ['当前阶段', status['current_phase']], ['总体进度', status['overall_progress_percent']]]
    progress_rows = [['ID', '名称', '时间', '状态', '进度', '核心讲解点']] + [[d['id'], d['name'], d['minutes'], d['status'], d['progress'], d['key_message']] for d in status['demos']]
    questions_rows = [['类型', '内容']] + [['用户问题', q] for q in status['user_questions']] + [['风险', r] for r in status['risks']] + [['下一步', a] for a in status['next_actions']]
    write_xlsx(out / 'course_dashboard.xlsx', {'Dashboard': dashboard_rows, 'Progress': progress_rows, 'Questions': questions_rows})

    build_html_gantt(status, out / 'gantt_dashboard.html')

    summary = f"""# Agent 生成摘要\n\n用户请求：{request_text.strip()}\n\n已生成：\n\n- {Path(actual_ppt).name}\n- course_dashboard.xlsx\n- gantt_dashboard.html\n\n脚本验证范围：文件存在、关键字段、基本结构。\n\n人工确认范围：课程进度真实性、时间安排合理性、用户问题完整性、风险表达是否合适。\n"""
    (out / 'agent_summary.md').write_text(summary, encoding='utf-8')
    manifest = {
        'generated_at': datetime.utcnow().isoformat() + 'Z',
        'inputs': {'status': args.status, 'request': args.request},
        'outputs': ['course_update.pptx', 'course_dashboard.xlsx', 'gantt_dashboard.html', 'agent_summary.md'],
        'validated_by_script': ['file existence', 'basic field presence'],
        'requires_human_review': ['content accuracy', 'teaching timing', 'business wording']
    }
    (out / 'context_manifest.json').write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding='utf-8')
    print('generated course assets in', out)

if __name__ == '__main__':
    main()
